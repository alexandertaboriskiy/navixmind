"""
Document Tools - PDF, Office documents (DOCX/PPTX/XLSX), file I/O, and ZIP archive creation
"""

import os
import zipfile
from io import BytesIO
from typing import Optional

from ..bridge import ToolError
from ..utils.file_limits import validate_file_for_processing, validate_pdf_for_processing, PROCESSING_LIMITS


def read_pdf(pdf_path: str, pages: str = "all") -> dict:
    """
    Extract text from a PDF file.

    Args:
        pdf_path: Path to the PDF file
        pages: Page range ("all", "1-5", "3", etc.)

    Returns:
        Dict with extracted text
    """
    from pypdf import PdfReader

    # Validate file
    validate_pdf_for_processing(pdf_path)

    try:
        reader = PdfReader(pdf_path)
        total_pages = len(reader.pages)

        # Parse page range
        if pages == "all":
            page_indices = range(total_pages)
        elif "-" in pages:
            start, end = pages.split("-")
            start = int(start) - 1
            end = min(int(end), total_pages)
            page_indices = range(start, end)
        else:
            page_num = int(pages) - 1
            if page_num >= total_pages:
                raise ToolError(f"Page {pages} doesn't exist. PDF has {total_pages} pages.")
            page_indices = [page_num]

        # Extract text
        text_parts = []
        for i in page_indices:
            page = reader.pages[i]
            text = page.extract_text()
            if text:
                text_parts.append(f"--- Page {i + 1} ---\n{text}")

        full_text = "\n\n".join(text_parts)

        # Truncate if too long
        if len(full_text) > 100000:
            full_text = full_text[:100000] + "\n\n[Content truncated...]"

        return {
            "path": pdf_path,
            "total_pages": total_pages,
            "pages_extracted": len(page_indices),
            "text": full_text
        }

    except Exception as e:
        raise ToolError(f"Failed to read PDF: {str(e)}")


def read_file(file_path: str) -> dict:
    """
    Read text content from a file.

    Args:
        file_path: Path to the file to read

    Returns:
        Dict with path, content, and size_bytes
    """
    validate_file_for_processing(file_path)

    try:
        size_bytes = os.path.getsize(file_path)
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()

        # Truncate if too long
        if len(content) > 100000:
            content = content[:100000] + "\n\n[Content truncated...]"

        return {
            "path": file_path,
            "content": content,
            "size_bytes": size_bytes,
        }

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to read file: {str(e)}")


def write_file(output_path: str, content: str) -> dict:
    """
    Write text content to a file.

    Args:
        output_path: Where to save the file
        content: Text content to write

    Returns:
        Dict with output_path, success, and size_bytes
    """
    max_chars = PROCESSING_LIMITS['text_chars']
    if len(content) > max_chars:
        raise ToolError(
            f"Content too large: {len(content)} chars. "
            f"Maximum: {max_chars} chars."
        )

    try:
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)

        size_bytes = os.path.getsize(output_path)

        return {
            "output_path": output_path,
            "success": True,
            "size_bytes": size_bytes,
        }

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to write file: {str(e)}")


def create_pdf(
    output_path: str,
    content: Optional[str] = None,
    title: Optional[str] = None,
    image_paths: Optional[list] = None,
) -> dict:
    """
    Create a PDF from text and/or images.

    Args:
        output_path: Where to save the PDF
        content: Optional text content for the PDF
        title: Optional document title
        image_paths: Optional list of image file paths to embed

    Returns:
        Dict with output path and page count
    """
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image

    if not content and not image_paths:
        raise ToolError("create_pdf requires at least 'content' or 'image_paths'")

    try:
        # Ensure output directory exists
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        # Create document
        page_width, page_height = letter
        margin = 72
        doc = SimpleDocTemplate(
            output_path,
            pagesize=letter,
            rightMargin=margin,
            leftMargin=margin,
            topMargin=margin,
            bottomMargin=margin,
        )

        usable_width = page_width - 2 * margin
        usable_height = page_height - 2 * margin

        # Get styles
        styles = getSampleStyleSheet()

        # Build content
        story = []

        # Add title if provided
        if title:
            title_style = ParagraphStyle(
                'Title',
                parent=styles['Heading1'],
                fontSize=18,
                spaceAfter=30,
            )
            story.append(Paragraph(title, title_style))
            story.append(Spacer(1, 0.25 * inch))

        # Add text content paragraphs
        if content:
            body_style = styles['Normal']
            for paragraph in content.split('\n\n'):
                if paragraph.strip():
                    # Escape special characters
                    safe_text = paragraph.replace('&', '&amp;')
                    safe_text = safe_text.replace('<', '&lt;')
                    safe_text = safe_text.replace('>', '&gt;')
                    story.append(Paragraph(safe_text, body_style))
                    story.append(Spacer(1, 0.1 * inch))

        # Embed images
        if image_paths:
            for img_path in image_paths:
                if not os.path.isfile(img_path):
                    raise ToolError(f"Image file not found: {img_path}")

                # Get image dimensions and scale to fit page
                from PIL import Image as PILImage
                with PILImage.open(img_path) as pil_img:
                    img_w, img_h = pil_img.size

                # Scale to fit within usable area while maintaining aspect ratio
                scale = min(usable_width / img_w, usable_height / img_h, 1.0)
                display_w = img_w * scale
                display_h = img_h * scale

                story.append(Image(img_path, width=display_w, height=display_h))
                story.append(Spacer(1, 0.2 * inch))

        # Build PDF
        doc.build(story)

        return {
            "output_path": output_path,
            "success": True,
        }

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to create PDF: {str(e)}")


def convert_document(
    input_path: str,
    output_format: str
) -> dict:
    """
    Convert document to another format.

    Args:
        input_path: Path to input document
        output_format: Target format (pdf, html, txt)

    Returns:
        Dict with output path
    """
    validate_file_for_processing(input_path, 'document')

    ext = os.path.splitext(input_path)[1].lower()
    base_name = os.path.splitext(input_path)[0]
    output_path = f"{base_name}.{output_format}"

    try:
        if ext in ['.docx', '.doc']:
            return _convert_docx(input_path, output_format, output_path)
        elif ext == '.txt':
            return _convert_txt(input_path, output_format, output_path)
        else:
            raise ToolError(f"Unsupported input format: {ext}")

    except Exception as e:
        raise ToolError(f"Conversion failed: {str(e)}")


def _convert_docx(input_path: str, output_format: str, output_path: str) -> dict:
    """Convert DOCX to target format."""
    from docx import Document

    doc = Document(input_path)

    # Extract text
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)

    text = '\n\n'.join(full_text)

    if output_format == 'txt':
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        return {"output_path": output_path, "success": True}

    elif output_format == 'pdf':
        return create_pdf(text, output_path)

    elif output_format == 'html':
        html = f"""<!DOCTYPE html>
<html>
<head><meta charset="utf-8"><title>Converted Document</title></head>
<body>
{''.join(f'<p>{p}</p>' for p in full_text if p.strip())}
</body>
</html>"""
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html)
        return {"output_path": output_path, "success": True}

    else:
        raise ToolError(f"Unsupported output format: {output_format}")


def _convert_txt(input_path: str, output_format: str, output_path: str) -> dict:
    """Convert TXT to target format."""
    with open(input_path, 'r', encoding='utf-8') as f:
        text = f.read()

    if output_format == 'pdf':
        return create_pdf(text, output_path)

    elif output_format == 'html':
        paragraphs = text.split('\n\n')
        html = f"""<!DOCTYPE html>
<html>
<head><meta charset="utf-8"><title>Converted Document</title></head>
<body>
{''.join(f'<p>{p}</p>' for p in paragraphs if p.strip())}
</body>
</html>"""
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html)
        return {"output_path": output_path, "success": True}

    else:
        raise ToolError(f"Unsupported output format: {output_format}")


def create_zip(
    output_path: str,
    file_paths: list,
    compression: str = "deflated",
) -> dict:
    """
    Create a ZIP archive from a list of files.

    Uses Python's built-in zipfile module (PSF License 2.0 — fully free/open).

    Args:
        output_path: Where to save the ZIP file
        file_paths: List of file paths to include in the archive
        compression: Compression method — "deflated" (smaller, default) or "stored" (no compression, faster)

    Returns:
        Dict with output path, file count, and total size
    """
    if not file_paths:
        raise ToolError("create_zip requires at least one file path")

    compression_map = {
        "deflated": zipfile.ZIP_DEFLATED,
        "stored": zipfile.ZIP_STORED,
    }
    if compression not in compression_map:
        raise ToolError(
            f"Unsupported compression: {compression}. Use 'deflated' or 'stored'."
        )

    try:
        # Ensure output directory exists
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        # Verify all files exist before creating the archive
        for fpath in file_paths:
            if not os.path.isfile(fpath):
                raise ToolError(f"File not found: {fpath}")

        # Track basenames to handle duplicates
        seen_names = {}
        zip_compression = compression_map[compression]

        with zipfile.ZipFile(output_path, 'w', compression=zip_compression) as zf:
            for fpath in file_paths:
                basename = os.path.basename(fpath)

                # Handle duplicate basenames by appending a counter
                if basename in seen_names:
                    seen_names[basename] += 1
                    name, ext = os.path.splitext(basename)
                    arcname = f"{name}_{seen_names[basename]}{ext}"
                else:
                    seen_names[basename] = 0
                    arcname = basename

                zf.write(fpath, arcname)

        # Get final archive size
        archive_size = os.path.getsize(output_path)

        return {
            "output_path": output_path,
            "success": True,
            "file_count": len(file_paths),
            "size_bytes": archive_size,
            "size_mb": round(archive_size / (1024 * 1024), 2),
            "compression": compression,
        }

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to create ZIP: {str(e)}")


# ---------------------------------------------------------------------------
# DOCX read/write
# ---------------------------------------------------------------------------

def read_docx(docx_path: str, extract: str = "all") -> dict:
    """
    Extract text, tables, and metadata from a DOCX file.

    Args:
        docx_path: Path to the DOCX file
        extract: What to extract — "text", "tables", or "all"

    Returns:
        Dict with extracted content
    """
    from docx import Document

    validate_file_for_processing(docx_path, 'document')

    try:
        doc = Document(docx_path)
        result = {"path": docx_path}

        if extract in ("text", "all"):
            paragraphs = [p.text for p in doc.paragraphs]
            full_text = "\n\n".join(p for p in paragraphs if p.strip())
            if len(full_text) > PROCESSING_LIMITS['text_chars']:
                full_text = full_text[:PROCESSING_LIMITS['text_chars']] + "\n\n[Content truncated...]"
            result["text"] = full_text
            result["paragraph_count"] = len(paragraphs)

        if extract in ("tables", "all"):
            tables_data = []
            for table in doc.tables:
                rows = []
                for row in table.rows:
                    rows.append([cell.text for cell in row.cells])
                tables_data.append(rows)
            result["tables"] = tables_data
            result["table_count"] = len(tables_data)

        return result

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to read DOCX: {str(e)}")


def modify_docx(input_path: str, output_path: str, operations: list) -> dict:
    """
    Modify an existing DOCX file.

    Args:
        input_path: Path to the source DOCX
        output_path: Where to save the modified DOCX
        operations: List of modifications, each with "action" and "params":
            - replace_text: {"old": str, "new": str}
            - add_paragraph: {"text": str, "style": str (optional)}
            - update_table_cell: {"table": int, "row": int, "col": int, "text": str}

    Returns:
        Dict with output path and operations applied
    """
    from docx import Document

    validate_file_for_processing(input_path, 'document')

    try:
        doc = Document(input_path)
        applied = 0

        for op in operations:
            action = op.get("action")
            params = op.get("params", {})

            if action == "replace_text":
                old = params.get("old", "")
                new = params.get("new", "")
                for para in doc.paragraphs:
                    if old in para.text:
                        for run in para.runs:
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                applied += 1

            elif action == "add_paragraph":
                text = params.get("text", "")
                style = params.get("style")
                doc.add_paragraph(text, style=style)
                applied += 1

            elif action == "update_table_cell":
                table_idx = params.get("table", 0)
                row_idx = params.get("row", 0)
                col_idx = params.get("col", 0)
                text = params.get("text", "")
                if table_idx < len(doc.tables):
                    table = doc.tables[table_idx]
                    if row_idx < len(table.rows) and col_idx < len(table.rows[row_idx].cells):
                        table.rows[row_idx].cells[col_idx].text = text
                        applied += 1

        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        doc.save(output_path)

        return {
            "output_path": output_path,
            "success": True,
            "operations_applied": applied,
        }

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to modify DOCX: {str(e)}")


# ---------------------------------------------------------------------------
# PPTX read/write
# ---------------------------------------------------------------------------

def read_pptx(pptx_path: str, extract: str = "all") -> dict:
    """
    Extract text, slide content, speaker notes, and metadata from a PPTX file.

    Args:
        pptx_path: Path to the PPTX file
        extract: What to extract — "text", "slides", "notes", or "all"

    Returns:
        Dict with extracted content
    """
    from pptx import Presentation

    validate_file_for_processing(pptx_path, 'document')

    try:
        prs = Presentation(pptx_path)

        if len(prs.slides) > PROCESSING_LIMITS.get('pptx_slides', 500):
            raise ToolError(
                f"Presentation has too many slides ({len(prs.slides)}). "
                f"Maximum: {PROCESSING_LIMITS.get('pptx_slides', 500)}"
            )

        result = {"path": pptx_path, "slide_count": len(prs.slides)}
        slides_data = []

        for i, slide in enumerate(prs.slides):
            slide_info = {"slide_number": i + 1, "shapes": []}

            # Extract text from shapes
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_text = "\n".join(p.text for p in shape.text_frame.paragraphs)
                    if shape_text.strip():
                        texts.append(shape_text)
                        slide_info["shapes"].append({
                            "name": shape.name,
                            "text": shape_text,
                        })

                if shape.has_table:
                    table_rows = []
                    for row in shape.table.rows:
                        table_rows.append([cell.text for cell in row.cells])
                    slide_info["shapes"].append({
                        "name": shape.name,
                        "table": table_rows,
                    })

            slide_info["text"] = "\n".join(texts)

            # Extract notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                slide_info["notes"] = slide.notes_slide.notes_text_frame.text
            else:
                slide_info["notes"] = ""

            slides_data.append(slide_info)

        if extract in ("text", "all"):
            all_text = []
            for sd in slides_data:
                all_text.append(f"--- Slide {sd['slide_number']} ---\n{sd['text']}")
            full_text = "\n\n".join(all_text)
            if len(full_text) > PROCESSING_LIMITS['text_chars']:
                full_text = full_text[:PROCESSING_LIMITS['text_chars']] + "\n\n[Content truncated...]"
            result["text"] = full_text

        if extract in ("slides", "all"):
            result["slides"] = slides_data

        if extract in ("notes", "all"):
            result["notes"] = [
                {"slide": sd["slide_number"], "notes": sd["notes"]}
                for sd in slides_data if sd["notes"]
            ]

        return result

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to read PPTX: {str(e)}")


def modify_pptx(input_path: str, output_path: str, operations: list) -> dict:
    """
    Modify an existing PPTX file.

    Args:
        input_path: Path to the source PPTX
        output_path: Where to save the modified PPTX
        operations: List of modifications, each with "action" and "params":
            - replace_text: {"old": str, "new": str}
            - add_slide: {"layout_index": int (default 1), "title": str, "content": str}
            - update_slide_text: {"slide": int (1-indexed), "shape_name": str, "text": str}
            - set_notes: {"slide": int (1-indexed), "text": str}

    Returns:
        Dict with output path and operations applied
    """
    from pptx import Presentation

    validate_file_for_processing(input_path, 'document')

    try:
        prs = Presentation(input_path)
        applied = 0

        for op in operations:
            action = op.get("action")
            params = op.get("params", {})

            if action == "replace_text":
                old = params.get("old", "")
                new = params.get("new", "")
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                for run in para.runs:
                                    if old in run.text:
                                        run.text = run.text.replace(old, new)
                applied += 1

            elif action == "add_slide":
                layout_idx = params.get("layout_index", 1)
                layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]
                slide = prs.slides.add_slide(layout)
                # Set title and content if placeholders exist
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 0 and params.get("title"):
                        ph.text = params["title"]
                    elif ph.placeholder_format.idx == 1 and params.get("content"):
                        ph.text = params["content"]
                applied += 1

            elif action == "update_slide_text":
                slide_num = params.get("slide", 1) - 1
                shape_name = params.get("shape_name", "")
                text = params.get("text", "")
                if 0 <= slide_num < len(prs.slides):
                    slide = prs.slides[slide_num]
                    for shape in slide.shapes:
                        if shape.name == shape_name and shape.has_text_frame:
                            shape.text_frame.paragraphs[0].text = text
                            applied += 1
                            break

            elif action == "set_notes":
                slide_num = params.get("slide", 1) - 1
                text = params.get("text", "")
                if 0 <= slide_num < len(prs.slides):
                    slide = prs.slides[slide_num]
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = text
                    applied += 1

        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        prs.save(output_path)

        return {
            "output_path": output_path,
            "success": True,
            "operations_applied": applied,
        }

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to modify PPTX: {str(e)}")


# ---------------------------------------------------------------------------
# XLSX read/write
# ---------------------------------------------------------------------------

def read_xlsx(xlsx_path: str, sheet: str = None, range: str = None, extract: str = "values") -> dict:
    """
    Extract cell data, sheet names, and formulas from an XLSX file.

    Args:
        xlsx_path: Path to the XLSX file
        sheet: Sheet name or index (omit for all sheets)
        range: Cell range like "A1:D10" (omit for all data)
        extract: What to extract — "values", "formulas", or "all"

    Returns:
        Dict with extracted data
    """
    from openpyxl import load_workbook

    validate_file_for_processing(xlsx_path, 'document')

    try:
        wb = load_workbook(xlsx_path, data_only=(extract == "values"))
        result = {
            "path": xlsx_path,
            "sheet_names": wb.sheetnames,
            "sheet_count": len(wb.sheetnames),
        }

        # Determine which sheets to process
        if sheet is not None:
            if sheet.isdigit():
                idx = int(sheet)
                if idx < len(wb.sheetnames):
                    sheets_to_read = [wb.sheetnames[idx]]
                else:
                    raise ToolError(f"Sheet index {idx} out of range. Available: {len(wb.sheetnames)} sheets.")
            elif sheet in wb.sheetnames:
                sheets_to_read = [sheet]
            else:
                raise ToolError(f"Sheet '{sheet}' not found. Available: {wb.sheetnames}")
        else:
            sheets_to_read = wb.sheetnames

        max_rows = PROCESSING_LIMITS.get('xlsx_rows', 100_000)
        sheets_data = {}

        for sheet_name in sheets_to_read:
            ws = wb[sheet_name]

            if range:
                cells = ws[range]
            else:
                cells = ws.iter_rows()

            rows_data = []
            row_count = 0
            for row in cells:
                if row_count >= max_rows:
                    rows_data.append(["[Truncated — max rows exceeded]"])
                    break
                if extract == "formulas":
                    rows_data.append([cell.value if cell.data_type == 'f' else cell.value for cell in row])
                else:
                    rows_data.append([cell.value for cell in row])
                row_count += 1

            sheets_data[sheet_name] = {
                "rows": rows_data,
                "row_count": row_count,
                "dimensions": ws.dimensions,
            }

        result["sheets"] = sheets_data
        wb.close()

        return result

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to read XLSX: {str(e)}")


def modify_xlsx(input_path: str, output_path: str, operations: list) -> dict:
    """
    Modify an existing XLSX file.

    Args:
        input_path: Path to the source XLSX
        output_path: Where to save the modified XLSX
        operations: List of modifications, each with "action" and "params":
            - set_cell: {"sheet": str, "cell": str, "value": any}
            - set_formula: {"sheet": str, "cell": str, "formula": str}
            - add_row: {"sheet": str, "values": list}
            - add_sheet: {"name": str}
            - delete_sheet: {"name": str}

    Returns:
        Dict with output path and operations applied
    """
    from openpyxl import load_workbook

    validate_file_for_processing(input_path, 'document')

    try:
        wb = load_workbook(input_path)
        applied = 0

        for op in operations:
            action = op.get("action")
            params = op.get("params", {})

            if action == "set_cell":
                sheet_name = params.get("sheet", wb.sheetnames[0])
                if sheet_name not in wb.sheetnames:
                    raise ToolError(f"Sheet '{sheet_name}' not found.")
                ws = wb[sheet_name]
                ws[params["cell"]] = params.get("value")
                applied += 1

            elif action == "set_formula":
                sheet_name = params.get("sheet", wb.sheetnames[0])
                if sheet_name not in wb.sheetnames:
                    raise ToolError(f"Sheet '{sheet_name}' not found.")
                ws = wb[sheet_name]
                ws[params["cell"]] = params.get("formula", "")
                applied += 1

            elif action == "add_row":
                sheet_name = params.get("sheet", wb.sheetnames[0])
                if sheet_name not in wb.sheetnames:
                    raise ToolError(f"Sheet '{sheet_name}' not found.")
                ws = wb[sheet_name]
                ws.append(params.get("values", []))
                applied += 1

            elif action == "add_sheet":
                name = params.get("name", "Sheet")
                wb.create_sheet(title=name)
                applied += 1

            elif action == "delete_sheet":
                name = params.get("name")
                if name in wb.sheetnames:
                    del wb[name]
                    applied += 1

        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        wb.save(output_path)
        wb.close()

        return {
            "output_path": output_path,
            "success": True,
            "operations_applied": applied,
        }

    except ToolError:
        raise
    except Exception as e:
        raise ToolError(f"Failed to modify XLSX: {str(e)}")
